from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[2]
OUT = ROOT / "docs/presentation/data-harbor-seidor-sales-deck-es.pptx"
LOGO = ROOT / "docs/presentation/assets/data-harbor-logo.png"
SEIDOR = ROOT / "docs/presentation/assets/seidor-logo.png"

BLUE = RGBColor(0, 112, 243)
NAVY = RGBColor(7, 13, 31)
MUTED = RGBColor(82, 94, 116)
LIGHT = RGBColor(246, 249, 252)
GREEN = RGBColor(19, 180, 115)
ORANGE = RGBColor(242, 164, 61)


SLIDES = [
    {
        "title": "DataHarbor by SEIDOR",
        "subtitle": "Calidad de datos de producto con agentes, evidencia y control humano.",
        "type": "Portada comercial",
        "message": "Tu catalogo puede moverse mas rapido sin perder control.",
        "image": None,
    },
    {
        "title": "La calidad del catalogo no escala con trabajo manual",
        "type": "Problema + coste oculto",
        "bullets": [
            "Cada producto trae atributos, formatos, categorias y reglas distintas.",
            "Cada campo incompleto retrasa publicacion, SEO, conversion y ventas.",
            "El coste real esta en buscar, validar, revisar y sincronizar una y otra vez.",
        ],
        "message": "El mantenimiento manual crea deuda operativa.",
        "image": None,
    },
    {
        "title": "Conecta tu catalogo y tus sistemas",
        "type": "Configuracion e integraciones",
        "bullets": [
            "Conecta Mirakl y sistemas PIM como Akeneo, Salsify, inriver o Sales Layer.",
            "Centraliza integraciones, investigacion, schemas, evidencia y export.",
            "Prepara el flujo para operar con sistemas reales, no con hojas sueltas.",
        ],
        "message": "Un punto de control para catalogo, fuentes y sincronizacion.",
        "image": "catalog-configuration.png",
    },
    {
        "title": "Importa productos y asigna schemas",
        "type": "Importacion y baseline",
        "bullets": [
            "Los productos importados se clasifican por schema y categoria.",
            "Se identifican productos que necesitan enriquecimiento.",
            "El equipo ve warnings y readiness antes de iniciar trabajo manual.",
        ],
        "message": "Una importacion se convierte en una cola priorizada de calidad.",
        "image": "catalog-baseline-schema-assignment.png",
    },
    {
        "title": "Visualiza la calidad del catalogo",
        "type": "Dashboard de calidad",
        "bullets": [
            "Vista general de productos, gaps, candidatos y evidencia.",
            "Filtros por necesidad de enriquecimiento, candidatos y warnings.",
            "Puntuacion de calidad para decidir donde actuar primero.",
        ],
        "message": "El equipo sabe que revisar y por que.",
        "image": "product-triage-dashboard.png",
    },
    {
        "title": "Detecta gaps antes de investigar",
        "type": "Comparacion inicial",
        "bullets": [
            "Compara el baseline Mirakl contra los campos candidatos.",
            "Muestra campos missing y warnings antes de ejecutar agentes.",
            "Prepara la siguiente accion desde la ficha de producto.",
        ],
        "message": "El sistema no solo reporta problemas; prepara la accion.",
        "image": "product-compare-empty-candidates.png",
    },
    {
        "title": "Lanza agentes para buscar datos externos",
        "type": "Agentes y cola de ejecucion",
        "bullets": [
            "Un agente busca informacion externa relevante para cada producto.",
            "Puede ejecutarse desde una ficha individual o en bloque para multiples productos.",
            "Cada ejecucion queda en cola con estado, runner, evidencia y candidatos.",
        ],
        "message": "Investigacion asistida, trazable y escalable.",
        "image": "research-agent-queue.png",
    },
    {
        "title": "Recibe datos estructurados con evidencia",
        "type": "Resultado del agente",
        "bullets": [
            "Compara valor actual, candidato y fuente.",
            "Mapea atributos como EAN, conectividad, peso, bateria y descripcion.",
            "Cada propuesta llega con evidencia para reducir retrabajo.",
        ],
        "message": "No es texto suelto: son cambios mapeados al producto.",
        "image": "product-data-comparison.png",
    },
    {
        "title": "Mejora descripcion, SEO y atributos criticos",
        "type": "SEO y contenido",
        "bullets": [
            "Identifica baja calidad SEO y descripciones incompletas.",
            "Propone descripciones enriquecidas y atributos tecnicos.",
            "Mantiene aprobacion humana antes de aplicar cambios.",
        ],
        "message": "Mejor contenido sin sacrificar control editorial.",
        "image": "candidate-attribute-review.png",
    },
    {
        "title": "Revisa, aprueba y sincroniza",
        "type": "Aprobacion humana",
        "bullets": [
            "Aprueba o rechaza candidatos por atributo.",
            "Conserva el baseline original y el candidato recomendado.",
            "Prepara cambios aprobados para sincronizar con tus sistemas.",
        ],
        "message": "Los agentes aceleran; el equipo decide.",
        "image": "candidate-attribute-review.png",
    },
    {
        "title": "Evidencia visible para cada decision",
        "type": "Trazabilidad de fuentes",
        "bullets": [
            "Diferencia fuentes de retailer, fabricante y otras referencias.",
            "Vincula evidencia al producto y a los campos que soporta.",
            "Permite abrir la fuente para revisar el origen.",
        ],
        "message": "La confianza viene de la fuente, no de una caja negra.",
        "image": "product-evidence-sources.png",
    },
    {
        "title": "Schemas por tipo de producto",
        "type": "Gobierno por categoria",
        "bullets": [
            "Configura familias de schema para categorias diferentes.",
            "Define campos requeridos, recomendados, warnings y scoring.",
            "Permite que la calidad sea especifica por categoria.",
        ],
        "message": "No todos los productos se revisan igual; DataHarbor lo entiende.",
        "image": "schema-configuration-overview.png",
    },
    {
        "title": "Reglas para calidad, SEO y scoring",
        "type": "Reglas operativas",
        "bullets": [
            "Campos requeridos y recomendados por schema.",
            "Warnings como brand missing, descripcion ausente o campo requerido faltante.",
            "Scoring rules para penalizar gaps de EAN, descripcion debil o informacion critica.",
        ],
        "message": "La calidad deja de depender de memoria tribal.",
        "image": "schema-field-requirements-and-rules.png",
    },
    {
        "title": "Aggregators con confianza y autoridad",
        "type": "Fuentes y confianza",
        "bullets": [
            "Define fuentes por autoridad: fabricante, retailer, base tecnica o referencia interna.",
            "Asigna confianza, cobertura y uso actual.",
            "Evita que fuentes de baja autoridad sobrescriban valores canonicos.",
        ],
        "message": "Gobierna de donde viene cada dato y cuanto confiar en el.",
        "image": "aggregator-configuration-overview.png",
    },
    {
        "title": "Configura reglas por fuente",
        "type": "Detalle de aggregator",
        "bullets": [
            "Define tipo de fuente, URL base, autoridad, confianza y dominios.",
            "Controla si una fuente participa en el workflow de revision.",
            "Ajusta la confianza por contexto y categoria.",
        ],
        "message": "Automatizacion configurable, auditable y segura.",
        "image": "official-manufacturer-aggregator-settings.png",
    },
    {
        "title": "Menos coste operativo. Mas calidad. Mas velocidad.",
        "type": "Impacto esperado",
        "bullets": [
            "Reduce horas manuales de busqueda y comparacion.",
            "Acelera publicacion de productos con datos completos.",
            "Mejora SEO y conversion con atributos y descripciones mas completas.",
            "Mantiene trazabilidad, governance y aprobacion humana.",
        ],
        "message": "DataHarbor convierte calidad de catalogo en ventaja operativa.",
        "image": "product-triage-dashboard.png",
    },
    {
        "title": "DataHarbor by SEIDOR",
        "subtitle": "Agentes que enriquecen. Evidencia que respalda. Equipos que aprueban.",
        "type": "Cierre comercial",
        "bullets": [
            "Conecta tus catalogos.",
            "Lanza agentes.",
            "Revisa evidencia.",
            "Sincroniza cambios con control.",
        ],
        "message": "La calidad del catalogo puede escalar sin perder control.",
        "image": None,
    },
]


def add_textbox(slide, x, y, w, h, text, size=24, color=NAVY, bold=False, align=None):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    if align:
        p.alignment = align
    return box


def add_header(slide, idx):
    slide.shapes.add_picture(str(LOGO), Inches(0.42), Inches(0.24), width=Inches(1.6))
    add_textbox(slide, 11.55, 0.28, 1.2, 0.28, f"{idx:02d}", 10, MUTED, align=PP_ALIGN.RIGHT)


def add_footer(slide, message):
    shape = slide.shapes.add_shape(1, Inches(0), Inches(6.78), Inches(13.333), Inches(0.72))
    shape.fill.solid()
    shape.fill.fore_color.rgb = NAVY
    shape.line.fill.background()
    add_textbox(slide, 0.55, 6.92, 11.9, 0.3, message, 13, RGBColor(255, 255, 255), bold=True)


def add_bullets(slide, bullets, x, y, w, h):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.size = Pt(15)
        p.font.color.rgb = NAVY
        p.space_after = Pt(9)
    return box


def add_image(slide, image_name):
    path = ROOT / "gamma" / image_name
    slide.shapes.add_shape(1, Inches(5.35), Inches(1.15), Inches(7.45), Inches(4.95)).fill.solid()
    frame = slide.shapes[-1]
    frame.fill.fore_color.rgb = RGBColor(255, 255, 255)
    frame.line.color.rgb = RGBColor(226, 232, 240)
    frame.line.width = Pt(1)
    slide.shapes.add_picture(str(path), Inches(5.52), Inches(1.32), width=Inches(7.1))


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    for idx, item in enumerate(SLIDES, start=1):
        slide = prs.slides.add_slide(blank)
        bg = slide.background.fill
        bg.solid()
        bg.fore_color.rgb = RGBColor(250, 252, 255)

        if idx in (1, 17):
            slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.333), Inches(7.5)).fill.solid()
            slide.shapes[-1].fill.fore_color.rgb = NAVY
            slide.shapes.add_picture(str(LOGO), Inches(0.75), Inches(0.68), width=Inches(2.3))
            slide.shapes.add_picture(str(SEIDOR), Inches(10.25), Inches(0.72), width=Inches(2.15))
            add_textbox(slide, 0.85, 2.25, 8.8, 0.8, item["title"], 44, RGBColor(255, 255, 255), bold=True)
            add_textbox(slide, 0.88, 3.15, 7.6, 0.7, item.get("subtitle", ""), 22, RGBColor(205, 220, 255))
            if item.get("bullets"):
                add_bullets(slide, item["bullets"], 0.92, 4.15, 4.4, 1.7)
            add_textbox(slide, 0.9, 6.58, 10.2, 0.3, item["message"], 16, RGBColor(255, 255, 255), bold=True)
            continue

        add_header(slide, idx)
        add_textbox(slide, 0.55, 0.9, 4.45, 0.22, item["type"].upper(), 8, BLUE, bold=True)
        add_textbox(slide, 0.55, 1.18, 4.65, 1.05, item["title"], 27, NAVY, bold=True)
        add_bullets(slide, item["bullets"], 0.65, 2.55, 4.2, 2.75)

        if item.get("image"):
            add_image(slide, item["image"])
        else:
            for i, (label, color) in enumerate([
                ("Coste manual", ORANGE),
                ("Datos incompletos", BLUE),
                ("Riesgo de publicacion", GREEN),
            ]):
                card = slide.shapes.add_shape(1, Inches(5.55 + i * 2.35), Inches(2.1), Inches(1.95), Inches(1.2))
                card.fill.solid()
                card.fill.fore_color.rgb = RGBColor(255, 255, 255)
                card.line.color.rgb = RGBColor(221, 228, 238)
                add_textbox(slide, 5.75 + i * 2.35, 2.45, 1.55, 0.38, label, 14, color, bold=True, align=PP_ALIGN.CENTER)

        add_footer(slide, item["message"])

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(OUT)


if __name__ == "__main__":
    build()
